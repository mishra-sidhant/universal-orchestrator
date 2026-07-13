from __future__ import annotations

from collections.abc import Sequence
from html import escape
import os
from pathlib import Path
import shutil
import subprocess
import tempfile
import textwrap
import zipfile

from universal_orchestrator.models import Artifact, ArtifactType, SlideSpec
from universal_orchestrator.utils import sha256_file


class ArtifactBuilder:
    def build_pdf(self, markdown: str, path: Path) -> Artifact:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet
        from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

        path.parent.mkdir(parents=True, exist_ok=True)
        styles = getSampleStyleSheet()
        story = []
        for line in markdown.splitlines():
            if not line.strip():
                story.append(Spacer(1, 8))
                continue
            if line.startswith("# "):
                style = styles["Heading1"]
            elif line.startswith("## "):
                style = styles["Heading2"]
            else:
                style = styles["BodyText"]
            text = line.lstrip("#").strip()
            story.append(Paragraph(escape(text), style))
            story.append(Spacer(1, 4))
        doc = SimpleDocTemplate(str(path), pagesize=A4, title="Universal Orchestrator Final Product")
        doc.build(story)
        return self._artifact(path, ArtifactType.PDF)

    def build_docx(self, markdown: str, path: Path) -> Artifact:
        from docx import Document

        path.parent.mkdir(parents=True, exist_ok=True)
        document = Document()
        for line in markdown.splitlines():
            if line.startswith("# "):
                document.add_heading(line[2:].strip(), level=1)
            elif line.startswith("## "):
                document.add_heading(line[3:].strip(), level=2)
            elif line.strip():
                document.add_paragraph(line.strip())
        document.save(str(path))
        return self._artifact(path, ArtifactType.DOCX)

    def validate_pdf(self, path: Path) -> list[str]:
        errors: list[str] = []
        try:
            from pypdf import PdfReader

            reader = PdfReader(str(path))
            if len(reader.pages) < 1:
                errors.append("PDF has no pages.")
        except Exception as exc:
            errors.append(f"PDF validation failed: {exc}")
        return errors

    def validate_docx(self, path: Path) -> list[str]:
        errors: list[str] = []
        try:
            from docx import Document

            document = Document(str(path))
            if not document.paragraphs:
                errors.append("DOCX has no paragraphs.")
        except Exception as exc:
            errors.append(f"DOCX validation failed: {exc}")
        return errors

    def build_pptx(self, slides: list[SlideSpec], path: Path) -> Artifact:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        path.parent.mkdir(parents=True, exist_ok=True)
        presentation = Presentation()
        for spec in slides:
            wrapped: list[str] = []
            for line in spec.body:
                wrapped.extend(textwrap.wrap(line, width=120) or [""])
            chunks = [wrapped[index:index + 5] for index in range(0, len(wrapped), 5)] or [[]]
            for chunk_index, body in enumerate(chunks, start=1):
                slide = presentation.slides.add_slide(presentation.slide_layouts[5])
                title = slide.shapes.title
                if title is not None:
                    title.text = spec.title if chunk_index == 1 else f"{spec.title} (continued)"
                textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.2))
                frame = textbox.text_frame
                frame.clear()
                for index, line in enumerate(body):
                    paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
                    paragraph.text = line
                    paragraph.font.size = Pt(20)
                if spec.notes:
                    slide.notes_slide.notes_text_frame.text = spec.notes
        presentation.save(str(path))
        return self._artifact(path, ArtifactType.PPTX)

    def validate_pptx(self, path: Path) -> list[str]:
        errors: list[str] = []
        try:
            from pptx import Presentation
            from pptx.util import Inches

            presentation = Presentation(str(path))
            slide_width = presentation.slide_width or 0
            slide_height = presentation.slide_height or 0
            if not presentation.slides:
                errors.append("PPTX has no slides.")
            for index, slide in enumerate(presentation.slides, start=1):
                if not any(getattr(shape, "text", "").strip() for shape in slide.shapes):
                    errors.append(f"PPTX slide {index} has no visible text.")
                for shape in slide.shapes:
                    left = shape.left or 0
                    top = shape.top or 0
                    width = shape.width or 0
                    height = shape.height or 0
                    if left < 0 or top < 0:
                        errors.append(f"PPTX slide {index} contains an off-canvas shape.")
                    if left + width > slide_width + Inches(0.05):
                        errors.append(f"PPTX slide {index} contains a horizontally clipped shape.")
                    if top + height > slide_height + Inches(0.05):
                        errors.append(f"PPTX slide {index} contains a vertically clipped shape.")
        except Exception as exc:
            errors.append(f"PPTX validation failed: {exc}")
        return errors

    def validate_rendered(
        self,
        kind: str,
        path: Path,
        quality_bar: str = "serious",
        expected_anchors: list[str] | None = None,
    ) -> tuple[list[str], list[str]]:
        """Run structural and bitmap-level checks without claiming visual perfection."""

        structural_validator = {
            "pdf": self.validate_pdf,
            "docx": self.validate_docx,
            "pptx": self.validate_pptx,
            "patch_plan": self.validate_patch_plan,
        }.get(kind)
        if structural_validator is None:
            return [f"No validator is registered for artifact kind {kind}."], []
        errors = structural_validator(path)
        if errors or kind == "patch_plan":
            return errors, []
        rendered: Path | None = None
        render_dir: Path | None = None
        try:
            rendered, render_error = self._render_first_page(kind, path)
            if render_error is not None:
                message = f"Render validation unavailable: {render_error}"
                return self._render_quality_result(message, quality_bar)
            if rendered is None:
                return self._render_quality_result(
                    "Render validation produced no preview image.", quality_bar
                )
            render_dir = rendered.parent
            page_paths = self._rendered_page_paths(rendered)
            expected_page_count = self._expected_rendered_page_count(kind, path)
            if expected_page_count is not None and len(page_paths) != expected_page_count:
                return self._render_quality_result(
                    f"Renderer produced {len(page_paths)} rendered pages; "
                    f"expected {expected_page_count} rendered pages.",
                    quality_bar,
                )
            page_errors = self._inspect_rendered_pages(page_paths, render_dir)
            if page_errors:
                return self._render_quality_result_list(page_errors, quality_bar)
            anchor_errors = self._validate_text_anchors(kind, path, expected_anchors or [])
            if anchor_errors:
                return anchor_errors, []
            return [], []
        finally:
            if render_dir is not None:
                shutil.rmtree(render_dir, ignore_errors=True)

    def _render_quality_result(self, message: str, quality_bar: str) -> tuple[list[str], list[str]]:
        if quality_bar in {"serious", "max"}:
            return [message], []
        return [], [message]

    def _render_quality_result_list(
        self, messages: list[str], quality_bar: str
    ) -> tuple[list[str], list[str]]:
        if quality_bar in {"serious", "max"}:
            return messages, []
        return [], messages

    def _rendered_page_paths(self, first_page: Path) -> list[Path]:
        pages = sorted(
            first_page.parent.glob("page-*.png"),
            key=lambda item: int(item.stem.rsplit("-", 1)[-1]),
        )
        return pages or [first_page]

    def _expected_rendered_page_count(self, kind: str, path: Path) -> int | None:
        try:
            if kind == "pdf":
                from pypdf import PdfReader

                return len(PdfReader(str(path)).pages)
            if kind == "pptx":
                from pptx import Presentation

                return len(Presentation(str(path)).slides)
        except Exception:
            return None
        return None

    def _inspect_rendered_pages(self, pages: list[Path], render_dir: Path) -> list[str]:
        from PIL import Image, ImageChops, ImageStat

        if not pages:
            return ["Rendered artifact produced no pages."]
        errors: list[str] = []
        thumbnails: list[Image.Image] = []
        for page_number, page in enumerate(pages, start=1):
            try:
                with Image.open(page) as image:
                    rgb = image.convert("RGB")
                    if rgb.width < 100 or rgb.height < 100:
                        errors.append(
                            f"Rendered artifact page {page_number} is unexpectedly small."
                        )
                    background = Image.new("RGB", rgb.size, "white")
                    difference = ImageChops.difference(rgb, background)
                    mask = difference.convert("L").point(
                        lambda value: 255 if value > 10 else 0
                    )
                    nonblank_ratio = ImageStat.Stat(mask).mean[0] / 255
                    if nonblank_ratio < 0.0005:
                        errors.append(
                            f"Rendered artifact page {page_number} is visually blank "
                            f"(nonblank ratio {nonblank_ratio:.6f})."
                        )
                    thumbnail = rgb.copy()
                    thumbnail.thumbnail((240, 180))
                    thumbnails.append(thumbnail)
            except Exception as exc:
                errors.append(
                    f"Rendered artifact page {page_number} could not be inspected: {exc}"
                )
        self._write_contact_sheet(thumbnails, render_dir)
        return errors

    def _write_contact_sheet(self, thumbnails: Sequence[object], render_dir: Path) -> None:
        from PIL import Image

        if not thumbnails:
            return
        typed_thumbnails = [item for item in thumbnails if isinstance(item, Image.Image)]
        if not typed_thumbnails:
            return
        columns = min(3, len(typed_thumbnails))
        rows = (len(typed_thumbnails) + columns - 1) // columns
        sheet = Image.new("RGB", (columns * 240, rows * 180), "white")
        for index, thumbnail in enumerate(typed_thumbnails):
            left = (index % columns) * 240
            top = (index // columns) * 180
            sheet.paste(thumbnail, (left, top))
        sheet.save(render_dir / "contact-sheet.png")

    def _validate_text_anchors(
        self, kind: str, path: Path, expected_anchors: list[str]
    ) -> list[str]:
        if not expected_anchors:
            return []
        try:
            if kind == "pdf":
                from pypdf import PdfReader

                text = "\n".join(page.extract_text() or "" for page in PdfReader(str(path)).pages)
            elif kind == "docx":
                from docx import Document

                text = "\n".join(paragraph.text for paragraph in Document(str(path)).paragraphs)
            elif kind == "pptx":
                from pptx import Presentation

                text = "\n".join(
                    getattr(shape, "text", "")
                    for slide in Presentation(str(path)).slides
                    for shape in slide.shapes
                )
            else:
                text = path.read_text(errors="replace")
        except Exception as exc:
            return [f"Artifact text anchors could not be inspected: {exc}"]
        return [
            f"Artifact is missing required text anchor: {anchor}"
            for anchor in expected_anchors
            if anchor not in text
        ]

    def _render_first_page(self, kind: str, path: Path) -> tuple[Path | None, str | None]:
        """Render every page and return the first page for compatibility with older callers."""

        pdftoppm = os.getenv("UO_PDFTOPPM_BIN") or shutil.which("pdftoppm")
        if not pdftoppm:
            return None, "pdftoppm is not installed"
        temp_dir = Path(tempfile.mkdtemp(prefix="uo-render-"))
        source_pdf = path
        render_succeeded = False
        try:
            if kind in {"docx", "pptx"}:
                soffice = os.getenv("UO_SOFFICE_BIN") or shutil.which("soffice")
                if not soffice:
                    return None, "LibreOffice soffice is not installed"
                completed = subprocess.run(
                    [
                        soffice,
                        "--headless",
                        "--convert-to",
                        "pdf",
                        "--outdir",
                        str(temp_dir),
                        str(path),
                    ],
                    capture_output=True,
                    text=True,
                    timeout=30,
                    check=False,
                )
                source_pdf = temp_dir / f"{path.stem}.pdf"
                if completed.returncode != 0 or not source_pdf.exists():
                    message = (
                        completed.stderr.strip()
                        or completed.stdout.strip()
                        or "conversion failed"
                    )
                    return None, f"{kind} conversion failed: {message}"
            output_prefix = temp_dir / "page"
            completed = subprocess.run(
                [pdftoppm, "-png", str(source_pdf), str(output_prefix)],
                capture_output=True,
                text=True,
                timeout=30,
                check=False,
            )
            pages = sorted(
                temp_dir.glob("page-*.png"),
                key=lambda item: int(item.stem.rsplit("-", 1)[-1]),
            )
            if completed.returncode != 0 or not pages:
                message = (
                    completed.stderr.strip()
                    or completed.stdout.strip()
                    or "PDF rasterization failed"
                )
                return None, message
            render_succeeded = True
            return pages[0], None
        except subprocess.TimeoutExpired as exc:
            return None, f"render subprocess timed out: {exc.cmd}"
        except OSError as exc:
            return None, f"render subprocess failed: {exc}"
        finally:
            if not render_succeeded:
                shutil.rmtree(temp_dir, ignore_errors=True)

    def build_patch_plan(self, markdown: str, path: Path) -> Artifact:
        path.parent.mkdir(parents=True, exist_ok=True)
        content = (
            "# Repository Patch Plan\n\n"
            "No repository changes were applied by this deterministic run. "
            "This file is a plan, not an implementation patch.\n\n"
            f"{markdown}"
        )
        path.write_text(content)
        return self._artifact(path, ArtifactType.REPORT)

    def validate_patch_plan(self, path: Path) -> list[str]:
        errors: list[str] = []
        if not path.exists():
            return ["Patch plan does not exist."]
        text = path.read_text(errors="replace")
        if not text.startswith("# Repository Patch Plan"):
            errors.append("Patch plan is missing its explicit plan heading.")
        if "not an implementation patch" not in text:
            errors.append("Patch plan does not disclose that no code patch was produced.")
        return errors

    def build_zip(self, artifacts: list[Artifact], path: Path) -> Artifact:
        path.parent.mkdir(parents=True, exist_ok=True)
        with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as archive:
            for artifact in sorted(artifacts, key=lambda item: item.name):
                artifact_path = artifact.as_path
                if artifact_path.exists() and artifact_path.resolve() != path.resolve():
                    archive.write(artifact_path, arcname=artifact_path.name)
        return self._artifact(path, ArtifactType.ZIP)

    def validate_zip(self, path: Path, expected_names: list[str] | None = None) -> list[str]:
        errors: list[str] = []
        if not path.exists():
            return ["ZIP file does not exist."]
        try:
            with zipfile.ZipFile(path) as archive:
                bad_member = archive.testzip()
                if bad_member:
                    errors.append(f"ZIP member failed CRC check: {bad_member}")
                names = archive.namelist()
                if not names:
                    errors.append("ZIP contains no files.")
                duplicates = sorted({name for name in names if names.count(name) > 1})
                errors.extend(f"ZIP contains duplicate member: {name}" for name in duplicates)
                if expected_names is not None:
                    expected = set(expected_names)
                    actual = set(names)
                    errors.extend(
                        f"ZIP is missing required member: {name}"
                        for name in sorted(expected - actual)
                    )
                    errors.extend(
                        f"ZIP contains unexpected member: {name}"
                        for name in sorted(actual - expected)
                    )
        except Exception as exc:
            errors.append(f"ZIP validation failed: {exc}")
        return errors

    def _artifact(self, path: Path, artifact_type: ArtifactType) -> Artifact:
        return Artifact(
            type=artifact_type,
            name=path.name,
            path=str(path),
            content_hash=sha256_file(path),
            size_bytes=path.stat().st_size,
        )
