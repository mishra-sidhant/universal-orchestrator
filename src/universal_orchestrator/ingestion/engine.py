from __future__ import annotations

import mimetypes
import csv
import io
import json
import tarfile
import urllib.request
import zipfile
from pathlib import Path
from urllib.parse import urlparse

from universal_orchestrator.ingestion.detectors import detect_input_type
from universal_orchestrator.ingestion.hardening import IngestionLimits, detect_text_encoding, symlink_warning
from universal_orchestrator.models import (
    ContextManifest,
    HostInvocation,
    InputAttachment,
    InputRecord,
    InputStatus,
    InputType,
    new_id,
)
from universal_orchestrator.repo import RepoAnalyzer
from universal_orchestrator.security import redact_text, scan_text
from universal_orchestrator.utils import compact_whitespace, iter_files, sha256_bytes, sha256_file, truncate_words


IGNORED_NAMES = {
    ".git",
    ".hg",
    ".svn",
    "node_modules",
    ".venv",
    "venv",
    "__pycache__",
    ".pytest_cache",
    ".mypy_cache",
    ".ruff_cache",
    "dist",
    "build",
}


class InputIngestor:
    def __init__(self, max_file_bytes: int = 5_000_000, max_folder_files: int = 500) -> None:
        self.limits = IngestionLimits(max_file_bytes=max_file_bytes, max_folder_files=max_folder_files)
        self.repo_analyzer = RepoAnalyzer()

    @property
    def max_file_bytes(self) -> int:
        return self.limits.max_file_bytes

    @property
    def max_folder_files(self) -> int:
        return self.limits.max_folder_files

    def ingest(self, invocation: HostInvocation, run_id: str) -> ContextManifest:
        records: list[InputRecord] = []
        warnings: list[str] = []

        prompt_findings = scan_text(invocation.prompt, location="prompt")
        prompt_text = redact_text(invocation.prompt)
        records.append(
            InputRecord(
                id="input_prompt",
                type=InputType.PROMPT,
                name="User prompt",
                uri="prompt://current",
                status=InputStatus.PARSED,
                content_hash=sha256_bytes(prompt_text.encode("utf-8")),
                summary=truncate_words(prompt_text, 120),
                content_text=prompt_text,
                metadata={"chars": len(prompt_text)},
                security_findings=prompt_findings,
            )
        )

        for attachment in invocation.attachments:
            try:
                records.append(
                    self._ingest_attachment(
                        attachment,
                        invocation.cwd,
                        invocation.user_options.allow_internet,
                    )
                )
            except Exception as exc:  # pragma: no cover - defensive boundary
                warnings.append(f"Failed to ingest {attachment.uri}: {exc}")
                records.append(
                    InputRecord(
                        id=new_id("input"),
                        type=detect_input_type(attachment.uri),
                        name=attachment.name or Path(attachment.uri).name or attachment.uri,
                        uri=attachment.uri,
                        status=InputStatus.FAILED,
                        warnings=[str(exc)],
                    )
                )

        for link in invocation.links:
            if not any(item.uri == link for item in records):
                records.append(self._ingest_url(InputAttachment(uri=link), invocation.cwd, invocation.user_options.allow_internet))

        return ContextManifest(
            run_id=run_id,
            invocation_id=invocation.id,
            prompt={
                "raw": prompt_text,
                "parsed_intent": self._infer_prompt_intent(invocation.prompt),
                "quality": invocation.user_options.quality,
                "budget_profile": invocation.user_options.budget_profile,
            },
            inputs=records,
            warnings=warnings,
        )

    def _ingest_attachment(
        self, attachment: InputAttachment, cwd: str | None, allow_network: bool
    ) -> InputRecord:
        input_type = detect_input_type(attachment.uri)
        if input_type in {InputType.URL, InputType.API}:
            return self._ingest_url(attachment, cwd, allow_network)

        original_path = Path(attachment.uri).expanduser()
        warning_for_symlink = symlink_warning(original_path)
        path = self._resolve_path(attachment.uri, cwd)
        name = attachment.name or path.name or str(path)
        if not path.exists():
            return InputRecord(
                id=new_id("input"),
                type=input_type,
                name=name,
                uri=attachment.uri,
                path=str(path),
                status=InputStatus.FAILED,
                warnings=[warning for warning in [warning_for_symlink, "Path does not exist."] if warning],
            )

        if input_type in {InputType.FOLDER, InputType.REPO}:
            record = self._ingest_folder(path, input_type, attachment)
            if warning_for_symlink:
                record.warnings.append(warning_for_symlink)
            return record
        if input_type in {InputType.TEXT, InputType.MARKDOWN, InputType.CODE, InputType.UNKNOWN}:
            record = self._ingest_text_file(path, input_type, attachment)
            if warning_for_symlink:
                record.warnings.append(warning_for_symlink)
            return record
        if input_type == InputType.PDF:
            return self._ingest_pdf(path, attachment)
        if input_type == InputType.DOCX:
            return self._ingest_docx(path, attachment)
        if input_type == InputType.PPTX:
            return self._ingest_pptx(path, attachment)
        if input_type == InputType.SPREADSHEET:
            return self._ingest_spreadsheet(path, attachment)
        if input_type == InputType.IMAGE:
            return self._ingest_image(path, attachment)
        if input_type == InputType.ARCHIVE:
            return self._ingest_archive(path, attachment)
        return self._ingest_binary_metadata(path, input_type, attachment)

    def _resolve_path(self, uri: str, cwd: str | None) -> Path:
        path = Path(uri).expanduser()
        if not path.is_absolute() and cwd:
            path = Path(cwd).expanduser() / path
        return path.resolve()

    def _ingest_text_file(
        self, path: Path, input_type: InputType, attachment: InputAttachment
    ) -> InputRecord:
        size = path.stat().st_size
        warnings: list[str] = []
        if size > self.max_file_bytes:
            warnings.append(f"File exceeds max_file_bytes={self.max_file_bytes}; summary uses prefix only.")
        raw = path.read_bytes()[: self.max_file_bytes]
        encoding = detect_text_encoding(raw)
        text = raw.decode(encoding, errors="replace")
        findings = scan_text(text, location=str(path))
        redacted = redact_text(text)
        return InputRecord(
            id=new_id("input"),
            type=input_type,
            name=attachment.name or path.name,
            uri=attachment.uri,
            path=str(path),
            status=InputStatus.PARSED,
            content_hash=sha256_file(path),
            size_bytes=size,
            mime_type=mimetypes.guess_type(path.name)[0],
            summary=truncate_words(redacted, 140),
            content_text=redacted,
            metadata={"chars_read": len(text), "suffix": path.suffix.lower(), "encoding": encoding},
            warnings=warnings,
            security_findings=findings,
        )

    def _ingest_docx(self, path: Path, attachment: InputAttachment) -> InputRecord:
        warnings: list[str] = []
        paragraphs: list[str] = []
        table_count = 0
        try:
            from docx import Document  # type: ignore[import-not-found]

            document = Document(path)
            paragraphs = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
            table_count = len(document.tables)
            for table in document.tables[:5]:
                for row in table.rows[:5]:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        paragraphs.append(" | ".join(cells))
        except Exception as exc:  # pragma: no cover - optional parser boundary
            warnings.append(f"DOCX parser failed: {exc}")
        text = "\n".join(paragraphs)
        findings = scan_text(text, location=str(path))
        redacted = redact_text(text)
        return InputRecord(
            id=new_id("input"),
            type=InputType.DOCX,
            name=attachment.name or path.name,
            uri=attachment.uri,
            path=str(path),
            status=InputStatus.PARSED if text else InputStatus.PARTIAL,
            content_hash=sha256_file(path),
            size_bytes=path.stat().st_size,
            mime_type=mimetypes.guess_type(path.name)[0],
            summary=truncate_words(redacted, 180)
            or "DOCX file parsed but no text was extracted.",
            content_text=redacted,
            metadata={"paragraphs": len(paragraphs), "tables": table_count, "suffix": ".docx"},
            warnings=warnings,
            security_findings=findings,
        )

    def _ingest_pptx(self, path: Path, attachment: InputAttachment) -> InputRecord:
        warnings: list[str] = []
        slide_text: list[str] = []
        slide_count = 0
        try:
            from pptx import Presentation  # type: ignore[import-not-found]

            presentation = Presentation(path)
            slide_count = len(presentation.slides)
            for index, slide in enumerate(presentation.slides[:50], start=1):
                fragments: list[str] = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        fragments.append(shape.text.strip())
                try:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        fragments.append(f"Notes: {notes}")
                except Exception:
                    pass
                if fragments:
                    slide_text.append(f"Slide {index}: " + " | ".join(fragments))
            if slide_count > 50:
                warnings.append("PPTX has more than 50 slides; MVP summary uses the first 50.")
        except Exception as exc:  # pragma: no cover - optional parser boundary
            warnings.append(f"PPTX parser failed: {exc}")
        text = "\n".join(slide_text)
        findings = scan_text(text, location=str(path))
        redacted = redact_text(text)
        return InputRecord(
            id=new_id("input"),
            type=InputType.PPTX,
            name=attachment.name or path.name,
            uri=attachment.uri,
            path=str(path),
            status=InputStatus.PARSED if text else InputStatus.PARTIAL,
            content_hash=sha256_file(path),
            size_bytes=path.stat().st_size,
            mime_type=mimetypes.guess_type(path.name)[0],
            summary=truncate_words(redacted, 180)
            or "PPTX file parsed but no slide text was extracted.",
            content_text=redacted,
            metadata={"slides": slide_count, "suffix": ".pptx"},
            warnings=warnings,
            security_findings=findings,
        )

    def _ingest_spreadsheet(self, path: Path, attachment: InputAttachment) -> InputRecord:
        suffix = path.suffix.lower()
        warnings: list[str] = []
        metadata: dict[str, object] = {"suffix": suffix}
        text_parts: list[str] = []
        if suffix in {".csv", ".tsv"}:
            delimiter = "\t" if suffix == ".tsv" else ","
            raw = path.read_text(errors="replace")
            reader = csv.reader(io.StringIO(raw), delimiter=delimiter)
            rows = []
            for index, row in enumerate(reader):
                rows.append(row)
                if index >= 24:
                    break
            metadata["sample_rows"] = len(rows)
            metadata["columns"] = max((len(row) for row in rows), default=0)
            text_parts = [" | ".join(row) for row in rows]
        else:
            try:
                import openpyxl  # type: ignore[import-not-found]

                workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
                metadata["sheets"] = workbook.sheetnames
                for sheet_name in workbook.sheetnames[:10]:
                    sheet = workbook[sheet_name]
                    rows = []
                    for index, row in enumerate(sheet.iter_rows(values_only=True)):
                        rows.append([str(cell) for cell in row if cell is not None])
                        if index >= 14:
                            break
                    text_parts.append(f"Sheet {sheet_name}: " + " / ".join(" | ".join(row) for row in rows if row))
                workbook.close()
            except Exception as exc:  # pragma: no cover - optional parser boundary
                warnings.append(f"Spreadsheet parser failed: {exc}")
        text = "\n".join(text_parts)
        findings = scan_text(text, location=str(path))
        redacted = redact_text(text)
        return InputRecord(
            id=new_id("input"),
            type=InputType.SPREADSHEET,
            name=attachment.name or path.name,
            uri=attachment.uri,
            path=str(path),
            status=InputStatus.PARSED if text else InputStatus.PARTIAL,
            content_hash=sha256_file(path),
            size_bytes=path.stat().st_size,
            mime_type=mimetypes.guess_type(path.name)[0],
            summary=truncate_words(redacted, 180)
            or "Spreadsheet parsed but no cell text was extracted.",
            content_text=redacted,
            metadata=metadata,
            warnings=warnings,
            security_findings=findings,
        )

    def _ingest_image(self, path: Path, attachment: InputAttachment) -> InputRecord:
        warnings: list[str] = []
        metadata: dict[str, object] = {"suffix": path.suffix.lower()}
        try:
            from PIL import Image  # type: ignore[import-not-found]

            with Image.open(path) as image:
                metadata.update(
                    {
                        "format": image.format,
                        "width": image.width,
                        "height": image.height,
                        "mode": image.mode,
                    }
                )
        except Exception as exc:
            warnings.append(f"Image metadata parser failed: {exc}")
        dimensions = (
            f"{metadata.get('width')}x{metadata.get('height')}"
            if metadata.get("width") and metadata.get("height")
            else "unknown dimensions"
        )
        return InputRecord(
            id=new_id("input"),
            type=InputType.IMAGE,
            name=attachment.name or path.name,
            uri=attachment.uri,
            path=str(path),
            status=InputStatus.PARSED if "width" in metadata else InputStatus.PARTIAL,
            content_hash=sha256_file(path),
            size_bytes=path.stat().st_size,
            mime_type=mimetypes.guess_type(path.name)[0],
            summary=f"Image metadata extracted: {dimensions}, format={metadata.get('format', 'unknown')}. OCR is not enabled yet.",
            content_text=f"Image metadata: dimensions={dimensions}; format={metadata.get('format', 'unknown')}.",
            metadata=metadata,
            warnings=warnings,
        )

    def _ingest_archive(self, path: Path, attachment: InputAttachment) -> InputRecord:
        warnings: list[str] = []
        metadata: dict[str, object] = {"suffix": path.suffix.lower(), "entries": 0, "unsafe_paths": []}
        entries: list[str] = []
        try:
            if zipfile.is_zipfile(path):
                with zipfile.ZipFile(path) as archive:
                    infos = archive.infolist()
                    metadata["entries"] = len(infos)
                    metadata["uncompressed_bytes"] = sum(info.file_size for info in infos)
                    entries = [info.filename for info in infos[:50]]
            elif tarfile.is_tarfile(path):
                with tarfile.open(path) as archive:
                    members = archive.getmembers()
                    metadata["entries"] = len(members)
                    metadata["uncompressed_bytes"] = sum(member.size for member in members)
                    entries = [member.name for member in members[:50]]
            else:
                warnings.append("Archive type is not supported for safe inventory yet.")
        except Exception as exc:
            warnings.append(f"Archive inventory failed: {exc}")
        unsafe = [name for name in entries if name.startswith("/") or ".." in Path(name).parts]
        metadata["unsafe_paths"] = unsafe
        if metadata.get("entries", 0) > self.limits.max_archive_entries:
            warnings.append(f"Archive exceeds max entries: {metadata['entries']}")
        if metadata.get("uncompressed_bytes", 0) > self.limits.max_archive_uncompressed_bytes:
            warnings.append(f"Archive exceeds max uncompressed bytes: {metadata['uncompressed_bytes']}")
        if unsafe:
            warnings.append(f"Archive contains unsafe paths: {unsafe[:5]}")
        summary = f"Archive inventory: {metadata.get('entries', 0)} entries. Sample: {', '.join(entries[:10])}"
        return InputRecord(
            id=new_id("input"),
            type=InputType.ARCHIVE,
            name=attachment.name or path.name,
            uri=attachment.uri,
            path=str(path),
            status=InputStatus.PARSED if entries or metadata.get("entries") == 0 else InputStatus.PARTIAL,
            content_hash=sha256_file(path),
            size_bytes=path.stat().st_size,
            mime_type=mimetypes.guess_type(path.name)[0],
            summary=summary,
            content_text=summary,
            metadata=metadata,
            warnings=warnings,
        )

    def _ingest_pdf(self, path: Path, attachment: InputAttachment) -> InputRecord:
        size = path.stat().st_size
        warnings: list[str] = []
        metadata: dict[str, object] = {"suffix": ".pdf"}
        text_parts: list[str] = []
        try:
            import pdfplumber  # type: ignore[import-not-found]

            with pdfplumber.open(path) as pdf:
                metadata["pages"] = len(pdf.pages)
                for index, page in enumerate(pdf.pages[:25], start=1):
                    page_text = page.extract_text(x_tolerance=1, y_tolerance=3) or ""
                    if page_text:
                        text_parts.append(f"Page {index}: {page_text}")
                if len(pdf.pages) > 25:
                    warnings.append("PDF has more than 25 pages; MVP summary uses the first 25 pages.")
        except Exception as exc:  # pragma: no cover - depends on optional PDF runtime
            warnings.append(f"PDF text extraction unavailable: {exc}")

        text = "\n".join(text_parts)
        findings = scan_text(text, location=str(path))
        redacted = redact_text(text)
        status = InputStatus.PARSED if text_parts else InputStatus.PARTIAL
        return InputRecord(
            id=new_id("input"),
            type=InputType.PDF,
            name=attachment.name or path.name,
            uri=attachment.uri,
            path=str(path),
            status=status,
            content_hash=sha256_file(path),
            size_bytes=size,
            mime_type="application/pdf",
            summary=truncate_words(redacted, 180),
            content_text=redacted,
            metadata=metadata,
            warnings=warnings,
            security_findings=findings,
        )

    def _ingest_binary_metadata(
        self, path: Path, input_type: InputType, attachment: InputAttachment
    ) -> InputRecord:
        size = path.stat().st_size
        summary_by_type = {
            InputType.DOCX: "DOCX file detected; structured parsing is planned for the next parser milestone.",
            InputType.PPTX: "PPTX file detected; slide parsing is planned for the next parser milestone.",
            InputType.SPREADSHEET: "Spreadsheet detected; schema extraction is planned for the next parser milestone.",
            InputType.IMAGE: "Image detected; OCR and visual source cards are planned for the next parser milestone.",
            InputType.ARCHIVE: "Archive detected; sandbox unpacking is planned and not performed in this MVP.",
            InputType.AUDIO_VIDEO: "Media file detected; transcription is planned for a later milestone.",
        }
        status = InputStatus.PARTIAL if input_type in summary_by_type else InputStatus.PARSED
        return InputRecord(
            id=new_id("input"),
            type=input_type,
            name=attachment.name or path.name,
            uri=attachment.uri,
            path=str(path),
            status=status,
            content_hash=sha256_file(path),
            size_bytes=size,
            mime_type=mimetypes.guess_type(path.name)[0],
            summary=summary_by_type.get(input_type, f"{input_type} file detected."),
            content_text=summary_by_type.get(input_type, f"{input_type} file detected."),
            metadata={"suffix": path.suffix.lower()},
            warnings=[] if status == InputStatus.PARSED else ["Structured parser not implemented yet."],
        )

    def _ingest_folder(
        self, path: Path, input_type: InputType, attachment: InputAttachment
    ) -> InputRecord:
        files = iter_files(path, IGNORED_NAMES, self.max_folder_files)
        suffix_counts: dict[str, int] = {}
        total_bytes = 0
        for file_path in files:
            suffix = file_path.suffix.lower() or "[none]"
            suffix_counts[suffix] = suffix_counts.get(suffix, 0) + 1
            try:
                total_bytes += file_path.stat().st_size
            except OSError:
                pass
        warnings = []
        if len(files) >= self.max_folder_files:
            warnings.append(f"Folder scan capped at {self.max_folder_files} files.")
        summary = (
            f"Scanned {len(files)} files under {path.name}. "
            f"Common suffixes: {self._format_top_counts(suffix_counts)}."
        )
        metadata = {
            "files_scanned": len(files),
            "total_scanned_bytes": total_bytes,
            "suffix_counts": suffix_counts,
            "is_git_repo": (path / ".git").exists(),
        }
        if input_type == InputType.REPO:
            metadata["repo_map"] = self.repo_analyzer.analyze(path, self.max_folder_files).model_dump(mode="json")
        return InputRecord(
            id=new_id("input"),
            type=input_type,
            name=attachment.name or path.name,
            uri=attachment.uri,
            path=str(path),
            status=InputStatus.PARSED,
            content_hash=sha256_bytes(compact_whitespace(str(metadata)).encode("utf-8")),
            size_bytes=total_bytes,
            summary=summary,
            content_text=summary,
            metadata=metadata,
            warnings=warnings,
        )

    def _ingest_url(self, attachment: InputAttachment, cwd: str | None, allow_network: bool) -> InputRecord:
        del cwd
        input_type = detect_input_type(attachment.uri)
        parsed = urlparse(attachment.uri)
        if allow_network:
            return self._fetch_url(attachment, input_type, parsed)
        return InputRecord(
            id=new_id("input"),
            type=input_type,
            name=attachment.name or parsed.netloc or attachment.uri,
            uri=attachment.uri,
            status=InputStatus.PARTIAL,
            content_hash=sha256_bytes(attachment.uri.encode("utf-8")),
            summary="URL recorded but not fetched because internet access requires explicit runtime permission.",
            content_text="URL recorded but not fetched because internet access requires explicit runtime permission.",
            metadata={"scheme": parsed.scheme, "netloc": parsed.netloc, "path": parsed.path},
            warnings=["URL fetch not performed in deterministic MVP."],
        )

    def _fetch_url(self, attachment: InputAttachment, input_type: InputType, parsed) -> InputRecord:
        warnings: list[str] = []
        metadata: dict[str, object] = {"scheme": parsed.scheme, "netloc": parsed.netloc, "path": parsed.path}
        text = ""
        status = InputStatus.PARTIAL
        try:
            request = urllib.request.Request(
                attachment.uri,
                headers={"User-Agent": "universal-orchestrator/0.1"},
                method="GET",
            )
            with urllib.request.urlopen(request, timeout=10) as response:
                content_type = response.headers.get("content-type")
                metadata["content_type"] = content_type
                raw = response.read(self.max_file_bytes)
                metadata["bytes_read"] = len(raw)
                text = raw.decode("utf-8", errors="replace")
                status = InputStatus.PARSED
                if input_type == InputType.API:
                    try:
                        parsed_json = json.loads(text)
                        metadata["json_type"] = type(parsed_json).__name__
                        text = json.dumps(parsed_json, indent=2)[: self.max_file_bytes]
                    except json.JSONDecodeError:
                        warnings.append("API URL did not return valid JSON.")
        except Exception as exc:
            warnings.append(f"URL fetch failed: {exc}")
        findings = scan_text(text, location=attachment.uri)
        redacted = redact_text(text)
        return InputRecord(
            id=new_id("input"),
            type=input_type,
            name=attachment.name or parsed.netloc or attachment.uri,
            uri=attachment.uri,
            status=status,
            content_hash=sha256_bytes((text or attachment.uri).encode("utf-8")),
            size_bytes=len(text.encode("utf-8")) if text else None,
            summary=truncate_words(redacted, 180) if text else "URL fetch attempted but no text was extracted.",
            content_text=redacted,
            metadata=metadata,
            warnings=warnings,
            security_findings=findings,
        )

    def _infer_prompt_intent(self, prompt: str) -> str:
        lowered = prompt.lower()
        if any(word in lowered for word in ["implement", "build", "code", "repo", "test"]):
            return "implementation"
        if any(word in lowered for word in ["report", "research", "pdf", "docx"]):
            return "research_and_artifact"
        if any(word in lowered for word in ["review", "audit", "critique"]):
            return "review_and_validation"
        return "orchestrated_task"

    def _format_top_counts(self, counts: dict[str, int]) -> str:
        if not counts:
            return "none"
        top = sorted(counts.items(), key=lambda item: item[1], reverse=True)[:5]
        return ", ".join(f"{suffix}={count}" for suffix, count in top)
